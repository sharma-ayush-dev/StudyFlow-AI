import logging
from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_pptx(path: str) -> dict:
    """
    Extracts text from a PPTX file, preserving slide order.
    """
    logger.info(f"Extracting PPTX: {path}")
    slides_text = []
    
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            slide_lines = [f"--- Slide {i+1} ---"]
            
            # Extract slide title if available
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_lines.append(f"Title: {title_text}")
            
            # Extract from all shapes
            for shape in slide.shapes:
                # Skip title as it was already handled
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            slide_lines.append(p_text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_vals:
                            slide_lines.append(" - ".join(row_vals))
            
            if len(slide_lines) > 1:
                slides_text.append("\n".join(slide_lines))
    except Exception as e:
        logger.error(f"Failed to extract PPTX: {e}")
        raise e
        
    extracted_text = "\n\n".join(slides_text)
    return {
        "text": extracted_text,
        "source_type": "pptx",
        "ocr_used": False,
        "metadata": {}
    }

def extract_ppt(path: str) -> dict:
    """
    Extracts text from a PPT file.
    First tries extract_pptx (in case it is actually a renamed pptx), then falls back to a binary/OLE string-scraping method.
    """
    import re
    logger.info(f"Extracting PPT: {path}")
    try:
        # Try pptx parser first in case of misnamed files
        res = extract_pptx(path)
        res["source_type"] = "ppt"
        return res
    except Exception as e:
        logger.warning(f"python-pptx failed for {path}, trying binary extraction fallback: {e}")
        try:
            with open(path, 'rb') as f:
                content = f.read()
            
            text_parts = []
            
            # Find ASCII printable strings (4 or more chars)
            ascii_strings = re.findall(b'[\x20-\x7e\r\n\t]{4,}', content)
            for s in ascii_strings:
                try:
                    decoded = s.decode('ascii').strip()
                    if decoded:
                        text_parts.append(decoded)
                except:
                    pass
            
            # Find UTF-16 strings (4 or more characters, i.e. 8 or more bytes)
            utf16_strings = re.findall(b'(?:[\x20-\x7e\r\n\t]\x00){4,}', content)
            for s in utf16_strings:
                try:
                    decoded = s.decode('utf-16le').strip()
                    if decoded:
                        text_parts.append(decoded)
                except:
                    pass
            
            # Combine and filter out garbage / metadata strings
            combined = "\n".join(text_parts)
            lines = []
            for line in combined.split('\n'):
                line = line.strip()
                if len(line) > 5 and not any(junk in line for junk in ['Current User', 'PowerPoint Document', 'Document Summary Information', 'SummaryInformation']):
                    lines.append(line)
            
            extracted_text = "\n".join(lines)
            return {
                "text": extracted_text,
                "source_type": "ppt",
                "ocr_used": False,
                "metadata": {"note": "Extracted via binary fallback"}
            }
        except Exception as ex:
            logger.error(f"Failed to extract PPT with fallback: {ex}")
            raise ex
